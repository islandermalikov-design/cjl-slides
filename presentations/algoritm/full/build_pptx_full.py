#!/usr/bin/env python3
"""Build the full 36-slide Algoritm deck as native PPTX."""
import math
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

PRIMARY = RGBColor(0x93, 0x40, 0x2B)
ACCENT = RGBColor(0x5A, 0x24, 0x17)
DEEPEST = RGBColor(0x3D, 0x16, 0x0E)
SECONDARY = RGBColor(0x14, 0x12, 0x12)
GREY_BG = RGBColor(0xF1, 0xEF, 0xEE)
STEEL = RGBColor(0xE7, 0xE9, 0xEA)
STEEL_DARK = RGBColor(0x38, 0x3D, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x14, 0x12, 0x12)
INK_MUTED = RGBColor(0x5A, 0x57, 0x57)
WHITE_MUTED = RGBColor(0xC9, 0xC4, 0xC4)
PINK = RGBColor(0xD9, 0xA8, 0x8F)
RULE = RGBColor(0xD8, 0xD4, 0xD3)
RULE_DARK = RGBColor(0x3A, 0x37, 0x37)
ROW_TINT = RGBColor(0xF3, 0xE9, 0xE6)

FONT = "Arial"
A = "assets/"
LOGO_COLOR = A + "logo.png"
LOGO_WHITE = A + "logo-white.png"
LOGO_ASPECT = 339 / 202

SW, SH = Inches(13.333), Inches(7.5)
ML, MR, MT = Inches(0.6), Inches(0.6), Inches(0.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_img_size_cache = {}


def img_size(path):
    if path not in _img_size_cache:
        with Image.open(path) as im:
            _img_size_cache[path] = im.size
    return _img_size_cache[path]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if color is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def add_shadow(shape, blur=90000, dist=26000, alpha=24000, color="1A1815", direction=5400000):
    spPr = shape._element.spPr
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    effect_lst = etree.SubElement(spPr, qn("a:effectLst"))
    outer = etree.SubElement(effect_lst, qn("a:outerShdw"))
    outer.set("blurRad", str(blur))
    outer.set("dist", str(dist))
    outer.set("dir", str(direction))
    outer.set("rotWithShape", "0")
    clr = etree.SubElement(outer, qn("a:srgbClr"))
    clr.set("val", color)
    a = etree.SubElement(clr, qn("a:alpha"))
    a.set("val", str(alpha))


def gradient_rect(s, l, t, w, h, c1, c2, angle=45):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    shp.line.fill.background()
    shp.fill.gradient()
    stops = shp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[-1].position = 1.0
    stops[-1].color.rgb = c2
    try:
        shp.fill.gradient_angle = angle
    except Exception:
        pass
    return shp


def connector(s, x1, y1, x2, y2, color, dashed=False, width=1.0):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.shadow.inherit = False
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dashed:
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return ln


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, first=False,
             space_after=0, spacing=1.0, font=FONT, letter_caps=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text.upper() if letter_caps else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return p


def picture_h(s, path, l, t, height):
    width = Emu(int(height * LOGO_ASPECT))
    return s.shapes.add_picture(path, l, t, height=height, width=width)


def photo_box(s, l, t, w, h, path, caption=None, cap_size=9, pad=Inches(0.07)):
    mat = rect(s, l, t, w, h, STEEL)
    add_shadow(mat)
    inner_l, inner_t = l + pad, t + pad
    inner_w, inner_h = w - Emu(int(pad * 2)), h - Emu(int(pad * 2))
    iw, ih = img_size(path)
    box_ar = inner_w / inner_h
    img_ar = iw / ih
    if img_ar > box_ar:
        dw = inner_w
        dh = Emu(int(inner_w / img_ar))
    else:
        dh = inner_h
        dw = Emu(int(inner_h * img_ar))
    dl = inner_l + Emu(int((inner_w - dw) / 2))
    dt = inner_t + Emu(int((inner_h - dh) / 2))
    s.shapes.add_picture(path, dl, dt, width=dw, height=dh)
    if caption:
        cap_h = Inches(0.3)
        cap_t = dt + dh - cap_h
        cb = rect(s, dl, cap_t, dw, cap_h, SECONDARY)
        tf = cb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = 0; tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        add_para(tf, caption, cap_size, WHITE, bold=True, first=True, spacing=1.05)
    return dl, dt, dw, dh


def data_table(s, l, t, w, headers, rows, col_widths=None, header_h=0.42, row_h=0.4, font_size=11):
    nrows = len(rows) + 1
    ncols = len(headers)
    total_h = Inches(header_h + row_h * len(rows))
    gframe = s.shapes.add_table(nrows, ncols, l, t, w, total_h)
    table = gframe.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    table.rows[0].height = Inches(header_h)
    for r in range(1, nrows):
        table.rows[r].height = Inches(row_h)

    def style_cell(cell, text, bg, color, bold, size):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.09); tf.margin_right = Inches(0.05)
        tf.margin_top = 0; tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = color

    for c, htext in enumerate(headers):
        style_cell(table.cell(0, c), htext.upper(), SECONDARY, WHITE, True, font_size - 1.5)
    for r, row in enumerate(rows, start=1):
        bg = ROW_TINT if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            style_cell(table.cell(r, c), str(val), bg, INK, False, font_size)
    return gframe


def tile_grid(s, l, t, w, h, items, cols, on_dark=False, gap=0.12):
    n = len(items)
    rows = math.ceil(n / cols)
    gp = Inches(gap)
    cw = Emu(int((w - gp * (cols - 1)) / cols))
    ch = Emu(int((h - gp * (rows - 1)) / rows)) if rows > 1 else h
    for i, item in enumerate(items):
        head, body = item if isinstance(item, tuple) else (item, None)
        r = i // cols
        c = i % cols
        x = l + Emu(int(c * (cw + gp)))
        y = t + Emu(int(r * (ch + gp)))
        bg = RGBColor(0x22, 0x20, 0x20) if on_dark else GREY_BG
        rect(s, x, y, cw, ch, bg)
        rect(s, x, y, Inches(0.03), ch, PRIMARY)
        _, tf = textbox(s, x + Inches(0.13), y, cw - Inches(0.24), ch, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, head, 11.5, WHITE if on_dark else INK, bold=True, first=True, spacing=1.08,
                 space_after=(3 if body else 0))
        if body:
            add_para(tf, body, 9.5, WHITE_MUTED if on_dark else INK_MUTED, spacing=1.25)
    return t + h


def feat_list(s, l, t, w, items, on_dark=False, row_h=0.95):
    y = t
    rh = Inches(row_h)
    for head, body in items:
        rect(s, l, y, Inches(0.035), rh - Inches(0.15), PRIMARY if not on_dark else PINK)
        _, tf = textbox(s, l + Inches(0.16), y, w - Inches(0.16), rh)
        add_para(tf, head, 12.5, PRIMARY if not on_dark else RGBColor(0xF0, 0x95, 0x8D),
                 bold=True, first=True, spacing=1.08, space_after=2)
        add_para(tf, body, 10.5, WHITE_MUTED if on_dark else INK_MUTED, spacing=1.28)
        y = y + rh
    return y


def stat_row(s, l, t, w, items, cols, on_dark=False, num_size=26, cap_size=9.5, row_h=1.35):
    n = len(items)
    gap = Inches(0.28)
    cw = Emu(int((w - gap * (cols - 1)) / cols))
    for i, (num, cap) in enumerate(items):
        r = i // cols
        c = i % cols
        x = l + Emu(int(c * (cw + gap)))
        y = t + Emu(int(r * Inches(row_h)))
        rect(s, x, y, cw, Pt(2.2), RULE_DARK if on_dark else RULE)
        _, tf = textbox(s, x, y + Inches(0.1), cw, Inches(0.5))
        add_para(tf, num, num_size, WHITE if on_dark else PRIMARY, bold=True, first=True)
        _, tf = textbox(s, x, y + Inches(0.62), cw, Inches(0.65))
        add_para(tf, cap, cap_size, WHITE_MUTED if on_dark else INK_MUTED, first=True, spacing=1.24)


def tag_row(s, l, t, tags, on_dark=False, max_w=None, size=10):
    x = l
    y = t
    h = Inches(0.36)
    for tag in tags:
        tw = Inches(0.28 + 0.085 * len(tag))
        if max_w is not None and (x + tw - l) > max_w:
            x = l
            y = y + h + Inches(0.12)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, tw, h)
        box.shadow.inherit = False
        box.fill.background()
        box.line.color.rgb = RGBColor(0x5A, 0x56, 0x56) if on_dark else RGBColor(0x8A, 0x86, 0x85)
        box.line.width = Pt(0.75)
        tf = box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        add_para(tf, tag, size, WHITE if on_dark else INK, bold=True, align=PP_ALIGN.CENTER, first=True)
        x = x + tw + Inches(0.12)
    return y + h


def icon_row(s, l, t, items, cols=8, gap=0.35, d=0.62):
    cw = Emu(int((Inches(12.1) - Inches(gap) * (cols - 1)) / cols))
    circle_d = Inches(d)
    for i, (glyph, label) in enumerate(items):
        x = l + Emu(int(i * (cw + Inches(gap))))
        cx = x + Emu(int((cw - circle_d) / 2))
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, t, circle_d, circle_d)
        circ.shadow.inherit = False
        circ.fill.solid(); circ.fill.fore_color.rgb = GREY_BG
        circ.line.fill.background()
        tf = circ.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_para(tf, glyph, 20, PRIMARY, bold=True, align=PP_ALIGN.CENTER, first=True)
        _, tf = textbox(s, x, t + circle_d + Inches(0.12), cw, Inches(0.5))
        add_para(tf, label, 9, INK, bold=True, align=PP_ALIGN.CENTER, first=True, spacing=1.15)


def kicker(s, l, t, w, text, color):
    _, tf = textbox(s, l, t, w, Inches(0.28))
    add_para(tf, text, 10.5, color, bold=True, first=True, letter_caps=True)


def page_head(s, kick, title, sub=None, on_dark=False, title_size=21, kick_color=None):
    kicker(s, ML, MT, Inches(9), kick, kick_color or (RGBColor(0xF0, 0xC3, 0xBF) if on_dark else PRIMARY))
    _, tf = textbox(s, ML, MT + Inches(0.32), Inches(11.5), Inches(0.9))
    add_para(tf, title, title_size, WHITE if on_dark else INK, bold=True, first=True, spacing=1.02)
    y = MT + Inches(0.32) + Inches(0.45 if len(title) < 55 else 0.8)
    if sub:
        _, tf = textbox(s, ML, y, Inches(11.5), Inches(0.55))
        add_para(tf, sub, 11.5, WHITE_MUTED if on_dark else INK_MUTED, first=True, spacing=1.25)
        y = y + Inches(0.55)
    return y + Inches(0.15)


_page_counter = [0]


def page_no(s, on_dark=False):
    _page_counter[0] += 1
    text = f"{_page_counter[0]:02d}"
    _, tf = textbox(s, Inches(11.9), Inches(6.95), Inches(1.1), Inches(0.35))
    add_para(tf, text, 10, WHITE_MUTED if on_dark else INK_MUTED, bold=True, align=PP_ALIGN.RIGHT, first=True)


def divider_slide(title_lines):
    s = slide()
    rect(s, 0, 0, SW, SH, STEEL_DARK)
    diamond = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(-3.0), Inches(9.5), Inches(9.5))
    diamond.rotation = 45
    diamond.shadow.inherit = False
    diamond.line.fill.background()
    diamond.fill.gradient()
    dstops = diamond.fill.gradient_stops
    dstops[0].position = 0.0
    dstops[0].color.rgb = PRIMARY
    dstops[-1].position = 1.0
    dstops[-1].color.rgb = STEEL_DARK
    try:
        diamond.fill.gradient_angle = 135
    except Exception:
        pass
    rect(s, Inches(0.6), Inches(3.35), Inches(1.0), Pt(3.2), PRIMARY)
    _, tf = textbox(s, Inches(0.6), Inches(3.6), Inches(8.5), Inches(2.2))
    for i, line in enumerate(title_lines):
        add_para(tf, line, 30, WHITE, bold=True, first=(i == 0), spacing=1.03)
    picture_h(s, LOGO_WHITE, Inches(11.05), Inches(6.98), Inches(0.32))
    return s


def gateway_diagram(s, bx, by, bw, bh, wireless, gateway_label, protocol_label, sensor_label="",
                     server_label="LAN", server_cap="Локальная сеть", n_sensors=3, gw_fill=None, srv_fill=None):
    line_color = PRIMARY if wireless else RGBColor(0x6E, 0x6A, 0x69)
    bh_half = Emu(int(bh / 2))
    mid_y = by + bh_half
    sensor_x = bx + Inches(0.28)
    sensor_r = Inches(0.07)
    if n_sensors == 3:
        sensor_ys = [by + Inches(0.12), mid_y, by + bh - Inches(0.12)]
    else:
        span = bh - Inches(0.24)
        sensor_ys = [by + Inches(0.12) + Emu(int(span * i / (n_sensors - 1))) for i in range(n_sensors)]

    gw_w, gw_h = Inches(1.15), Inches(0.46)
    gw_x = bx + Inches(1.0)
    gw_y = mid_y - Emu(int(gw_h / 2))

    srv_w, srv_h = Inches(1.15), Inches(0.46)
    srv_x = bx + bw - srv_w - Inches(0.12)
    srv_y = mid_y - Emu(int(srv_h / 2))

    for sy in sensor_ys:
        connector(s, sensor_x + sensor_r, sy, gw_x, mid_y, line_color, dashed=wireless, width=1.1)
    for sy in sensor_ys:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, sensor_x, sy - sensor_r, sensor_r * 2, sensor_r * 2)
        dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = SECONDARY
        dot.line.fill.background()

    if sensor_label:
        _, tf = textbox(s, sensor_x - Inches(0.4), by + bh - Inches(0.01), Inches(1.1), Inches(0.2))
        add_para(tf, sensor_label, 6.2, INK_MUTED, align=PP_ALIGN.CENTER, first=True)

    _, tf = textbox(s, gw_x - Inches(0.25), gw_y - Inches(0.2), gw_w + Inches(0.5), Inches(0.18))
    add_para(tf, protocol_label, 6.2, INK_MUTED, align=PP_ALIGN.CENTER, first=True)

    gbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gw_x, gw_y, gw_w, gw_h)
    gbox.shadow.inherit = False
    gbox.fill.solid(); gbox.fill.fore_color.rgb = gw_fill or (line_color if wireless else SECONDARY)
    gbox.line.fill.background()
    try:
        gbox.adjustments[0] = 0.18
    except Exception:
        pass
    tf = gbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(tf, gateway_label, 9.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)

    connector(s, gw_x + gw_w, mid_y, srv_x, mid_y, srv_fill or (PRIMARY if not wireless else SECONDARY), width=1.1)

    sbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, srv_x, srv_y, srv_w, srv_h)
    sbox.shadow.inherit = False
    sbox.fill.solid(); sbox.fill.fore_color.rgb = srv_fill or (PRIMARY if not wireless else SECONDARY)
    sbox.line.fill.background()
    try:
        sbox.adjustments[0] = 0.18
    except Exception:
        pass
    tf = sbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(tf, server_label, 9, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, srv_x - Inches(0.3), srv_y + srv_h + Inches(0.02), srv_w + Inches(0.6), Inches(0.18))
    add_para(tf, server_cap, 6.2, INK_MUTED, align=PP_ALIGN.CENTER, first=True)


def bar_chart(s, x, y, w, h, title_lines, categories, values, bar_color, first_color, cat_color, val_color):
    _, tf = textbox(s, x, y, w, Inches(0.5))
    add_para(tf, title_lines[0], 11.5, val_color, bold=True, first=True, spacing=1.15, space_after=1)
    for extra in title_lines[1:]:
        add_para(tf, extra, 11.5, val_color, bold=True, spacing=1.15)

    title_h = Inches(0.15) * len(title_lines) + Inches(0.35)
    top = y + title_h
    y_base = y + h - Inches(0.32)
    avail_h = y_base - top - Inches(0.28)
    max_val = max(values)

    connector(s, x, y_base, x + w, y_base, RGBColor(0x8A, 0x86, 0x85), width=0.75)

    n = len(values)
    slot_w = Emu(int(w / n))
    bar_w = Emu(int(slot_w * 0.5))
    for i, (cat, val) in enumerate(zip(categories, values)):
        bar_h = Emu(int(avail_h * (val / max_val)))
        slot_x = x + Emu(slot_w * i)
        bar_x = slot_x + Emu(int((slot_w - bar_w) / 2))
        bar_y = y_base - bar_h
        color = first_color if i == 0 else bar_color
        rect(s, bar_x, bar_y, bar_w, bar_h, color)
        _, tf = textbox(s, bar_x - Inches(0.15), bar_y - Inches(0.28), bar_w + Inches(0.3), Inches(0.24))
        add_para(tf, f"{val}%", 11, color, bold=True, align=PP_ALIGN.CENTER, first=True)
        _, tf = textbox(s, slot_x, y_base + Inches(0.05), slot_w, Inches(0.3))
        add_para(tf, cat, 8, cat_color, bold=True, align=PP_ALIGN.CENTER, first=True, spacing=1.1)


# ============================================================ SLIDE 1 — COVER
s = slide()
rect(s, 0, 0, SW, SH, STEEL)
gradient_rect(s, Inches(6.67), 0, Inches(6.67), SH, PRIMARY, ACCENT, angle=115)

picture_h(s, LOGO_COLOR, Inches(0.6), Inches(0.5), Inches(0.42))
_, tf = textbox(s, Inches(1.35), Inches(0.46), Inches(1.75), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "Комплексный подход к решению задач", 8, INK_MUTED, bold=True, first=True, spacing=1.15)

_, tf = textbox(s, Inches(0.6), Inches(2.55), Inches(5.5), Inches(2.3))
add_para(tf, "Промышленный мониторинг и предиктивное", 27, INK, bold=True, first=True, spacing=1.02)
add_para(tf, "обслуживание оборудования ГОК", 27, INK, bold=True, spacing=1.02)

_, tf = textbox(s, Inches(0.6), Inches(5.0), Inches(2.5), Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "Коммерческое предложение", 12.5, INK_MUTED, bold=True, first=True)
rect(s, Inches(2.95), Inches(5.0), Inches(0.72), Inches(0.42), PRIMARY)
_, tf = textbox(s, Inches(2.95), Inches(5.0), Inches(0.72), Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "2026", 12.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)

# monitor mockup with SuperCare screenshot
mon_w = Inches(4.55)
mon_pad = Inches(0.1)
img_w = mon_w - Emu(int(mon_pad * 2))
iw, ih = img_size(A + "supercare_left.jpg")
img_h = Emu(int(img_w * ih / iw))
mon_h = img_h + Emu(int(mon_pad * 2))
mon_l = Inches(8.15)
mon_t = Inches(1.85)
frame = rect(s, mon_l, mon_t, mon_w, mon_h, RGBColor(0x1B, 0x1B, 0x1B))
add_shadow(frame, blur=110000, dist=30000, alpha=32000)
s.shapes.add_picture(A + "supercare_left.jpg", mon_l + mon_pad, mon_t + mon_pad, width=img_w, height=img_h)
stand_w, stand_h = Inches(0.7), Inches(0.22)
rect(s, mon_l + Emu(int((mon_w - stand_w) / 2)), mon_t + mon_h, stand_w, stand_h, RGBColor(0x1B, 0x1B, 0x1B))
base_w, base_h = Inches(1.5), Inches(0.06)
rect(s, mon_l + Emu(int((mon_w - base_w) / 2)), mon_t + mon_h + stand_h, base_w, base_h, RGBColor(0x1B, 0x1B, 0x1B))

_, tf = textbox(s, Inches(0.6), Inches(6.95), Inches(6.0), Inches(0.35))
add_para(tf, "Инженерно-сервисный центр: г. Красноярск, 2026.", 10.5, INK_MUTED, first=True)
page_no(s)

# ============================================================ SLIDE 2 — ACHIEVEMENTS (moved to front)
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Результаты внедрения", "Измеримые результаты предиктивного обслуживания",
              "Совокупный эффект от внедрения технологии на промышленных объектах — снижение затрат "
              "на техническое обслуживание оценивается в 4 000 000 $.", title_size=21)
hero_stats = [
    ("−57%", "отказов механического оборудования"),
    ("−80%", "времени простоя из-за механических отказов"),
    ("−37%", "затрат на аутсорсинг технического обслуживания"),
]
hw = Inches(3.85)
gap = Inches(0.3)
for i, (num, cap) in enumerate(hero_stats):
    x = ML + Emu(int(i * (hw + gap)))
    rect(s, x, y + Inches(0.15), hw, Pt(2.8), PRIMARY)
    _, tf = textbox(s, x, y + Inches(0.3), hw, Inches(1.1))
    add_para(tf, num, 44, PRIMARY, bold=True, first=True)
    _, tf = textbox(s, x, y + Inches(1.35), hw, Inches(0.6))
    add_para(tf, cap, 13.5, INK, bold=True, first=True, spacing=1.25)
    _, tf = textbox(s, x, y + Inches(2.05), hw, Inches(0.9))
    add_para(tf, "Данные подтверждены на основе анализа 34 000+ успешных кейсов в горнодобывающей отрасли.",
             9, RGBColor(0x8A, 0x86, 0x85), first=True, spacing=1.35)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 3 — ABOUT
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(2.4), STEEL)
kicker(s, ML, Inches(0.4), Inches(5), "О компании", PRIMARY)
_, tf = textbox(s, ML, Inches(0.76), Inches(10.5), Inches(1.4))
add_para(tf, "ООО «Алгоритм» — ваш партнер в цифровизации", 22, INK, bold=True, first=True, spacing=1.02)
add_para(tf, "промышленности", 22, INK, bold=True, spacing=1.02)
cols = [
    (ML, None, "Мы предлагаем полный цикл внедрения систем прогнозного технического обслуживания: "
     "от аудита оборудования и подбора аппаратных решений до развертывания программного "
     "обеспечения и обучения персонала.", SECONDARY),
    (Inches(4.05), "Индивидуальный подход",
     "Адаптируем решения под специфику конкретного предприятия и отрасли.", PRIMARY),
    (Inches(7.15), "Экспертная поддержка",
     "Обеспечиваем консалтинг и сопровождение на всех этапах жизненного цикла системы.", PRIMARY),
    (Inches(10.25), "Мировые стандарты",
     "Опираемся на проверенную технологическую базу, успешно зарекомендовавшую себя на "
     "крупнейших промышленных объектах.", PRIMARY),
]
col_w = Inches(2.85)
top_y = Inches(2.85)
for l, heading, body, rule_color in cols:
    rect(s, l, top_y, col_w, Pt(2.5), rule_color)
    y = top_y + Inches(0.18)
    _, tf = textbox(s, l, y, col_w, Inches(3.2))
    first = True
    if heading:
        add_para(tf, heading, 14, PRIMARY, bold=True, first=True, spacing=1.1, space_after=6)
        first = False
    add_para(tf, body, 12, INK_MUTED if heading else INK, first=first, spacing=1.32)
picture_h(s, LOGO_COLOR, ML, Inches(7.0), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 4 — SCALE & TECH BASE (merged 3+4+26)
s = slide()
rect(s, 0, 0, SW, SH, STEEL)
y = page_head(s, "Технологическая база", "Масштаб и технологическая база платформы", title_size=21)
connector(s, Inches(6.67), y, Inches(6.67), Inches(7.0), RULE, width=0.75)

col_x = [ML, Inches(6.95)]
col_headings = ["Масштаб внедрения", "Научный потенциал платформы"]
col_stats = [
    [("35+", "стран — география работы технологической платформы на крупнейших промышленных объектах мира"),
     ("188 000+", "единиц оборудования находится под непрерывным онлайн-мониторингом платформы"),
     ("882 000+", "успешно установленных и работающих датчиков вибрации и температуры"),
     ("3.0 ТБ+", "данных телеметрии обрабатывается интеллектуальными ИИ-алгоритмами ежедневно")],
    [("42%", "команды — доля специалистов НИОКР (R&D) в общей численности персонала платформы"),
     ("250+", "патентов и зарегистрированных прав на собственное ПО и аппаратные решения"),
     ("800+", "сотрудников — общий штат экосистемы платформы, включая команду инженеров-диагностов"),
     ("2", "национальных стандарта, разработанных на базе технологии техобслуживания оборудования")],
]
for cx, heading, stats in zip(col_x, col_headings, col_stats):
    kicker(s, cx, y, Inches(5.5), heading, PRIMARY)
    yy = y + Inches(0.4)
    for num, cap in stats:
        rect(s, cx, yy, Inches(5.6), Pt(2), RULE)
        _, tf = textbox(s, cx, yy + Inches(0.1), Inches(5.6), Inches(0.6))
        add_para(tf, num, 22, INK, bold=True, first=True)
        _, tf = textbox(s, cx, yy + Inches(0.62), Inches(5.6), Inches(0.55))
        add_para(tf, cap, 10, INK_MUTED, first=True, spacing=1.2)
        yy = yy + Inches(1.15)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 5 — APPLICATIONS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Область применения", "Диагностируемое вращающееся оборудование",
              "Цель — обеспечить долгосрочную надёжную и безопасную эксплуатацию оборудования, "
              "избежать незапланированных простоев и вторичных аварий.", title_size=19)
tile_grid(s, ML, y + Inches(0.1), Inches(12.1), Inches(3.2), [
    ("Мельницы", "Шаровые и полуавтоматические мельницы"),
    ("Компрессоры", "Винтовые и воздушные компрессоры"),
    ("Насосы", "Шламовые и центробежные насосы"),
    ("Конвейеры", "Ленточные конвейеры и приводы"),
    ("Подъёмные машины", "Шахтные подъёмные установки"),
    ("Вентиляторы", "Вентиляторы главного проветривания"),
    ("Дробилки", "Щековые, конусные, роторные дробилки"),
    ("Мостовые краны", "Грузоподъёмное оборудование"),
    ("Прокатные станы", "Металлургическое оборудование"),
    ("Прочее", "Ж/д транспорт, ЦБП, винодельческое и медицинское оборудование"),
], cols=5)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

prs.save("algoritm-full-partial.pptx")
print("saved batch 1 (slides 1-5)")


# ============================================================ SLIDE 6 — ARCHITECTURE
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
page_head(s, "Архитектура", "Архитектура системы: надёжная передача данных и гибкая интеграция",
          title_size=19)


def arch_tier(s, y, label, cards, card_w_list=None, row_h=0.98):
    kicker(s, ML, y, Inches(11.5), label, PRIMARY)
    y2 = y + Inches(0.3)
    n = len(cards)
    gap = Inches(0.18)
    total_w = Inches(12.1)
    if card_w_list is None:
        cw = Emu(int((total_w - gap * (n - 1)) / n))
        card_w_list = [cw] * n
    x = ML
    row_h = Inches(row_h)
    for (head, body), cw in zip(cards, card_w_list):
        rect(s, x, y2, cw, row_h, GREY_BG)
        rect(s, x, y2, Inches(0.035), row_h, PRIMARY)
        _, tf = textbox(s, x + Inches(0.16), y2 + Inches(0.1), cw - Inches(0.3), row_h - Inches(0.2))
        add_para(tf, head, 11.5, INK, bold=True, first=True, spacing=1.12, space_after=2)
        add_para(tf, body, 9.3, INK_MUTED, spacing=1.24)
        x = x + cw + gap
    return y2 + row_h


y = Inches(1.75)
y = arch_tier(s, y, "ВЕРХНИЙ УРОВЕНЬ — КУДА СОБИРАЮТСЯ ДАННЫЕ", [
    ("Локальный сервер предприятия (On-Premise) / облачный сервер",
     "Обеспечивает 100% безопасность данных внутри закрытого периметра завода или гибкий доступ через облако."),
    ("Программный комплекс SuperCare",
     "Интеллектуальный центр анализа больших данных и ИИ-диагностики."),
])
_, tf = textbox(s, Inches(6.15), y, Inches(1.0), Inches(0.3))
add_para(tf, "↓", 16, RGBColor(0xAF, 0xAB, 0xAA), first=True)
y = y + Inches(0.32)
y = arch_tier(s, y, "ЦЕНТРАЛЬНЫЙ УРОВЕНЬ — ИНФРАСТРУКТУРА", [
    ("Интеграция в локальную сеть предприятия (LAN / ВОЛС)",
     "Передача данных на сервер по защищённым промышленным протоколам."),
], card_w_list=[Inches(12.1)])
_, tf = textbox(s, Inches(6.15), y, Inches(1.0), Inches(0.3))
add_para(tf, "↓", 16, RGBColor(0xAF, 0xAB, 0xAA), first=True)
y = y + Inches(0.32)
y = arch_tier(s, y, "НИЖНИЙ УРОВЕНЬ — ТРИ ВАРИАНТА СБОРА ДАННЫХ С АГРЕГАТОВ", [
    ("Вариант 1. Высокоскоростной проводной сбор (кабель / оптика / Wi-Fi)",
     "Для оборудования: критически важные и энергонасыщенные агрегаты с постоянным питанием — "
     "дробилки, главные приводы, масляные прессы."),
    ("Вариант 2. Беспроводная сеть малой дальности (ZigBee)",
     "Для оборудования: компактные узлы и конвейерные линии на средних дистанциях — "
     "шахтные подъёмники, скребки, сушильные барабаны."),
    ("Вариант 3. Дальнобойная беспроводная сеть (LoRa / 4G / 5G)",
     "Для оборудования: удалённые объекты в радиусе до 1 км без прокладки кабельных трасс — "
     "воздуходувки, насосные станции, компрессоры."),
], row_h=1.15)

picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 7 — WIRELESS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Сценарии применения", "Беспроводной вариант решения", title_size=21)
data_table(s, ML, y + Inches(0.1), Inches(5.6),
           ["Оборудование", "Датчиков", "Тип датчика", "Тип шлюза"],
           [["Консольный насос", "4", "RH505", "RH560"],
            ["Двухопорный насос", "6", "RH605", "RH560"],
            ["Вентилятор", "6", "RH506", "RH570"]],
           col_widths=[1.9, 1.0, 1.3, 1.4])
_, tf = textbox(s, ML, y + Inches(1.75), Inches(5.6), Inches(0.9))
add_para(tf, "Одноосные и трёхосные беспроводные датчики вибрации (RH505/RH605, RW506/RW606) "
             "передают данные на шлюз по радиоканалу — без прокладки кабельных трасс.",
         11, INK_MUTED, first=True, spacing=1.3)
feat_list(s, ML, y + Inches(2.55), Inches(5.6), [
    ("🛠️ Монтаж за 15 минут",
     "Датчики фиксируются на оборудовании с помощью шпилек, магнитов или промышленного клея. "
     "Никаких кабельных трасс и штробления."),
    ("🔋 До 3 лет автономной работы",
     "Датчики оснащены встроенными промышленными батареями повышенной ёмкости. Замена элементов "
     "питания не требует демонтажа самого датчика."),
    ("📡 Помехозащищённый радиоканал",
     "Данные передаются на промышленный шлюз в автоматическом режиме. Система устойчива к "
     "электромагнитным помехам от мощных электродвигателей."),
], row_h=0.92)
photo_box(s, Inches(6.55), y + Inches(0.1), Inches(3.0), Inches(2.0), A + "wireless_components.jpg")
photo_box(s, Inches(9.65), y + Inches(0.1), Inches(3.05), Inches(2.0), A + "wireless_software.jpg")
photo_box(s, Inches(6.55), y + Inches(2.2), Inches(6.15), Inches(1.85), A + "wireless_install.jpg",
          "Монтаж беспроводных датчиков на промышленном объекте")
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 8 — WIRED
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Сценарии применения", "Проводной вариант решения", title_size=21)
data_table(s, ML, y + Inches(0.1), Inches(5.6),
           ["Оборудование", "Датчиков", "Тип датчика", "Тип шлюза"],
           [["Гориз. дробилка", "10", "RH103T", "RH2000"],
            ["Верт. дробилка", "16", "RH104T", "RH2000"],
            ["Зубч. дробилка", "13", "RH113T", "RH2000"]],
           col_widths=[1.9, 1.0, 1.3, 1.4])
_, tf = textbox(s, ML, y + Inches(1.75), Inches(5.6), Inches(0.9))
add_para(tf, "Проводные датчики с верхним или боковым (360°) выходом кабеля подключаются к шлюзу "
             "RH2000 — стационарная схема для ответственных агрегатов с постоянным питанием.",
         11, INK_MUTED, first=True, spacing=1.3)
feat_list(s, ML, y + Inches(2.55), Inches(5.6), [
    ("⚡ Постоянное питание и непрерывный сбор",
     "Схема с подключением к шлюзу RH2000 обеспечивает ежесекундный мониторинг без задержек. "
     "Идеально для агрегатов класса «А», где секунда простоя стоит миллионы рублей."),
    ("🛡️ Защита кабеля на 360°",
     "Датчики поставляются с верхним или боковым выходом кабеля в металлической бронерукаве. "
     "Полная защита от механических повреждений, пыли, влаги и вибрации."),
    ("🎯 Анализ переходных процессов",
     "Проводная схема позволяет ловить дефекты в моменты пуска и останова агрегата, когда нагрузка "
     "на подшипники и редукторы максимальна."),
], row_h=0.98)
photo_box(s, Inches(6.55), y + Inches(0.1), Inches(3.0), Inches(2.0), A + "wired_sensors.jpg")
photo_box(s, Inches(9.65), y + Inches(0.1), Inches(3.05), Inches(2.0), A + "wired_gateway.jpg")
photo_box(s, Inches(6.55), y + Inches(2.2), Inches(6.15), Inches(1.45), A + "wired_software.jpg")
_, tf = textbox(s, Inches(6.55), y + Inches(3.7), Inches(6.15), Inches(0.55))
add_para(tf, "Программный комплекс обеспечивает глубокий спектральный анализ формы волны вибрации в "
             "реальном времени, выявляя скрытые трещины и микросколы задолго до их визуального проявления.",
         9, INK_MUTED, first=True, spacing=1.25)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 9 — CASE STORY
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Кейс из практики", "От аномалии до устранения дефекта — за 5 часов",
              "Объект: подземный резервный погрузчик (критически важное оборудование шахты)", title_size=20)

steps = [
    (A + "case_chart1.jpg", "🛑 02:00 ночи", "Внезапный скачок вибрации",
     "Телеметрия в реальном времени зафиксировала резкое аномальное усиление вибрации на двигателе "
     "погрузчика. Система автоматически отправила сигнал тревоги."),
    (A + "case_control_room.jpg", "🧠 +15 минут", "Мгновенный экспертный анализ",
     "Дежурный эксперт удалённого диагностического центра проанализировал спектр волны и выдал "
     "заключение: «Критический дефект подшипника. Риск заклинивания двигателя»."),
    (A + "case_bearing.jpg", "🛠️ 06:43 утра", "Точечный плановый ремонт",
     "Бригада заехала в шахту и оперативно заменила повреждённый узел. Сепаратор подшипника был "
     "полностью разрушен — до заклинивания и пожара оставались считанные часы."),
    (A + "case_chart2.jpg", "✅ 07:00 утра", "Возврат в штатный режим",
     "После замены подшипника форма волны и графики вибрации мгновенно вернулись к норме. "
     "Погрузчик запущен в работу."),
]
cw = Inches(2.85)
gap = Inches(0.2)
photo_h = Inches(1.5)
for i, (img, time_lbl, head, body) in enumerate(steps):
    x = ML + Emu(int(i * (cw + gap)))
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.34), Inches(0.34))
    badge.shadow.inherit = False
    badge.fill.solid(); badge.fill.fore_color.rgb = PRIMARY
    badge.line.fill.background()
    tf = badge.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(tf, str(i + 1), 12, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, x + Inches(0.44), y + Inches(0.03), cw - Inches(0.44), Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, time_lbl, 11, SECONDARY, bold=True, first=True)
    photo_box(s, x, y + Inches(0.44), cw, photo_h, img)
    _, tf = textbox(s, x, y + Inches(0.44) + photo_h + Inches(0.1), cw, Inches(0.32))
    add_para(tf, head, 12.5, INK, bold=True, first=True, spacing=1.1)
    _, tf = textbox(s, x, y + Inches(0.44) + photo_h + Inches(0.42), cw, Inches(1.1))
    add_para(tf, body, 9.5, INK_MUTED, first=True, spacing=1.28)
    if i < 3:
        _, tf = textbox(s, x + cw + Inches(0.01), y + Inches(0.44) + Emu(int(photo_h / 2)) - Inches(0.15),
                         gap - Inches(0.02), Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, "→", 14, RGBColor(0xAF, 0xAB, 0xAA), bold=True, align=PP_ALIGN.CENTER, first=True)

result_y = Inches(6.05)
plaque = rect(s, ML, result_y, Inches(12.1), Inches(0.68), GREY_BG)
rect(s, ML, result_y, Inches(0.05), Inches(0.68), PRIMARY)
_, tf = textbox(s, ML + Inches(0.25), result_y + Inches(0.03), Inches(0.7), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "📈", 20, PRIMARY, first=True)
_, tf = textbox(s, ML + Inches(1.0), result_y + Inches(0.06), Inches(10.8), Inches(0.58))
add_para(tf, "Главный вывод: результат для заказчика", 11, PRIMARY, bold=True, first=True, space_after=1)
add_para(tf, "Предотвращено внезапное аварийное отключение, исключены миллионные экономические потери "
             "от простоя линии и ликвидирована угроза нарушению техники безопасности (ТБ) под землёй.",
         9, INK, spacing=1.15)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

prs.save("algoritm-full-partial.pptx")
print("saved batch 2 (slides 6-9)")


# ============================================================ SLIDE 11 — INDUSTRY SCENARIOS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Область применения",
              "Сценарии применения: контроль оборудования на всех этапах — от добычи до плавки",
              "Сокращаем внеплановые простои на 30–40%, окупаемость — 11–14 месяцев", title_size=19)
stat_row(s, ML, y + Inches(0.05), Inches(12.1), [
    ("−30–40%", "простоев оборудования"),
    ("11 мес.", "средняя окупаемость проекта"),
    ("до 50 ед.", "оборудования на участке под контролем"),
], cols=3, num_size=26, row_h=1.05)
tile_grid(s, ML, y + Inches(1.35), Inches(12.1), Inches(1.35), [
    ("Добыча полезных ископаемых", "Шахтные подъёмные машины, ленточные конвейеры, буровые установки, вентиляторы главного проветривания"),
    ("Переработка полезных ископаемых", "Щёковые и конусные дробилки, шаровые мельницы, флотационное и обогатительное оборудование"),
    ("Плавка и металлургия", "Прокатные станы, винтовые компрессоры, насосные станции, мостовые краны"),
], cols=3)
icon_row(s, ML, y + Inches(2.9), [
    ("⚙", "Мельница"), ("▽", "Дробилка"), ("○", "Насос"), ("✳", "Вентилятор"),
    ("▭", "Конвейер"), ("⌐", "Подъёмная машина"), ("▤", "Компрессор"), ("⊓", "Мостовой кран"),
])
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)


def class_scenario_slide(title, wireless, gw_label, protocol, sensor_label, srv_label, srv_cap,
                          photo_key, photo_note, features, gw_fill=None, srv_fill=None, n_sensors=3,
                          sub=None, feat_row_h=0.98):
    s = slide()
    rect(s, 0, 0, SW, SH, WHITE)
    y = page_head(s, "Сценарии применения", title, sub, title_size=18)
    diagram_box_l, diagram_box_t = ML, y + Inches(0.1)
    diagram_box_w, diagram_box_h = Inches(6.85), Inches(2.5)
    rect(s, diagram_box_l, diagram_box_t, diagram_box_w, diagram_box_h, GREY_BG)
    gateway_diagram(s, diagram_box_l + Inches(0.3), diagram_box_t + Inches(0.25), diagram_box_w - Inches(0.6),
                     diagram_box_h - Inches(0.5), wireless, gw_label, protocol, sensor_label,
                     srv_label, srv_cap, n_sensors=n_sensors, gw_fill=gw_fill, srv_fill=srv_fill)
    photo_box(s, ML, diagram_box_t + diagram_box_h + Inches(0.2), Inches(3.35), Inches(1.5), A + photo_key)
    rect(s, ML + Inches(3.5), diagram_box_t + diagram_box_h + Inches(0.2), Inches(3.35), Inches(1.5), GREY_BG)
    _, tf = textbox(s, ML + Inches(3.65), diagram_box_t + diagram_box_h + Inches(0.35), Inches(3.05), Inches(1.2),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, photo_note, 10.5, INK_MUTED, first=True, spacing=1.3)
    feat_list(s, Inches(7.75), y + Inches(0.1), Inches(4.95), features, row_h=feat_row_h)
    page_no(s)
    return s


class_scenario_slide(
    "Мониторинг вспомогательного оборудования (Класс C/D): автоматизация ручных инспекций",
    True, "Шлюз LoRa", "LoRa, до 1 км", "Датчики RW161", "LAN", "Локальная сеть",
    "cd_install_photos.jpg", "RW161 — беспроводные датчики вибрации версии LoRa для оборудования уровня C/D",
    [
        ("🛰 Дальнобойность до 1 км без кабеля",
         "Радиус действия канала достигает 1 км в условиях промышленной застройки. Один шлюз собирает данные со всей фабрики без прокладки кабельных трасс."),
        ("👥 Сокращение ручного труда на 50%",
         "Полная замена ежедневных плановых обходов персонала. Система сама непрерывно мониторит узлы и вызывает инженера только по сигналу тревоги."),
        ("🤖 Умная сигнализация аварий",
         "Интеллектуальный анализ автоматически настраивает базовый порог вибрации для каждого узла и мгновенно оповещает о зарождении дефекта."),
        ("📊 Сбор формы волны (уникальная функция)",
         "В отличие от стандартных датчиков других брендов, RW161 умеет собирать и передавать полную форму волны для детального экспертного анализа."),
        ("🌋 Взрывозащита промышленного уровня",
         "Корпус имеет максимальную степень защиты IP68 / ATEX. Датчики сертифицированы для работы в самых тяжёлых, влажных и взрывоопасных условиях."),
    ], n_sensors=3,
    sub="Беспроводные датчики вибрации RW161 разработаны специально для массового контроля вспомогательных узлов. "
        "Они позволяют перевести обслуживание сотен агрегатов в автоматический режим с минимальными затратами.",
    feat_row_h=0.86)

class_scenario_slide(
    "Мониторинг основного оборудования (Класс «B»): глубокий анализ неисправностей",
    True, "RH570", "ZigBee", "RH505/605, RW625", "LAN", "Локальная сеть",
    "b_fault_photos.jpg", "RH570 — беспроводная станция сбора данных, датчики по ZigBee",
    [
        ("🎯 Диагностика «3 в 1» (вибрация, температура, скорость)",
         "Датчик одновременно измеряет три критических параметра в широком диапазоне частот (до 20 кГц). Это позволяет обнаружить малейшие зарождающиеся дефекты подшипников и шестерён."),
        ("🔍 Сверхглубокий сбор данных (до 819 200 линий)",
         "Длина записи сигнала до 2048 Кб даёт возможность инженерам-диагностам проводить расширенный спектральный анализ и безошибочно определять тип поломки."),
        ("🛡 Фильтрация ложных сигналов",
         "Интеллектуальные алгоритмы умеют автоматически определять моменты пуска/останова и отсекать фоновые шумы соседних машин, исключая ложные срабатывания тревоги."),
        ("🔋 До 3 лет автономной работы",
         "Несмотря на высокую плотность сбора данных, датчики оптимизированы по энергопотреблению и работают до 3 лет на одной батарее в режиме циклического мониторинга."),
        ("🌋 Максимальная взрывозащита IP68",
         "Оборудование сертифицировано для жёстких условий эксплуатации. Корпус полностью герметичен, защищён от агрессивной химии, пыли и влаги."),
    ], n_sensors=4,
    sub="Высокопроизводительные беспроводные станции и датчики обеспечивают детальный спектральный анализ. "
        "Это идеальное решение для средних и крупных агрегатов, требующих точной локализации дефектов без прокладки кабеля.",
    feat_row_h=0.86)

class_scenario_slide(
    "Мониторинг критического оборудования (Класс «А»): непрерывный контроль высокой точности",
    False, "RH2000 / Exd", "RS485 / WiFi / 4-20mA", "IEPE, обороты, частицы масла", "SuperCare",
    "Система мониторинга",
    "a_fault_photos.jpg", "RH2000 / RH2000 Exd — взрывозащищённая станция сбора данных",
    [
        ("⛓ Комплексный многофакторный мониторинг",
         "Полная совместимость с датчиками любого типа: IEPE-датчики вибрации, лазерные датчики оборотов, а также системы контроля металлических частиц в масле."),
        ("⏱ Непрерывный сбор данных без пауз",
         "Запись параметров каждую секунду (индикация тренда — 30 с, детальная форма волны — каждые 30 мин). Система не пропустит дефект даже при мгновенных скачках нагрузки."),
        ("🧠 Диагностика экспертного уровня с помощью ИИ",
         "Анализ данных по более чем 50 специальным индексам. Искусственный интеллект автоматически сопоставляет вибрацию, обороты и качество масла для выявления скрытых угроз."),
        ("💥 Взрывозащищённое исполнение корпуса (Exd)",
         "Станция поставляется в специальном бронированном корпусе RH2000. Система полностью сертифицирована для работы во взрывоопасных зонах шахт и обогатительных фабрик."),
    ], n_sensors=4,
    sub="Стационарные взрывозащищённые комплексы RH2000 обеспечивают непрерывный посекундный сбор данных. "
        "Это бескомпромиссное решение для ключевых агрегатов предприятия, где риски внезапного останова недопустимы.",
    feat_row_h=1.0)

# ============================================================ SLIDE 15 — FAULT CONTROL
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Диагностика", "Матрица диагностики: автоматический контроль и классификация неисправностей",
              title_size=18)
fl_h = 1.0
feat_list(s, ML, y + Inches(0.05), Inches(5.6), [
    ("🔩 Дефекты привода", "Дисбаланс, несоосность, ослабление, структурный резонанс, эксцентриситет ротора, изгиб вала"),
    ("⚙ Дефекты подшипников", "Сколы, износ, раковины, ржавчина, трещины, деформация и поломка сепаратора, недостаток смазки"),
    ("🔄 Неисправности редуктора", "Недостаток смазки, плохое зацепление, износ и поломка зубьев, усталостное разрушение"),
    ("🛢 Проблемы со смазкой", "Загрязнение масла, химический износ, попадание посторонних предметов"),
], row_h=fl_h)
photo_box(s, Inches(6.55), y + Inches(0.05), Inches(6.15), Inches(4 * fl_h), A + "defect_photo_grid.jpg")
page_no(s)

# ============================================================ SLIDE 17 — SUPERCARE
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Программное обеспечение", "Программный комплекс SuperCare: экосистема предиктивного анализа верхнего уровня",
              "Единая цифровая платформа для централизованного мониторинга, глубокой ИИ-диагностики неисправностей "
              "и оперативного управления надёжностью всего парка оборудования холдинга.", title_size=18)
photo_box(s, ML, y + Inches(0.1), Inches(5.9), Inches(2.05), A + "supercare_left.jpg")
photo_box(s, Inches(6.75), y + Inches(0.1), Inches(5.95), Inches(2.05), A + "supercare_right.jpg")
tile_grid(s, ML, y + Inches(2.3), Inches(12.1), Inches(1.9), [
    ("💻 Доступ из любого браузера (B/S архитектура)",
     "Система не требует установки программ на рабочие места. Доступ к панели мониторинга защищён и возможен с любого компьютера в локальной сети завода."),
    ("🔌 Бесшовная совместимость с АСУ ТП",
     "Лёгкая интеграция с существующими на предприятии системами автоматизации (DCS, PLC, MES, SCADA). Данные аналитики встраиваются в цифровой контур фабрики."),
    ("🤖 Обработка больших данных (Big Data)",
     "Высокопроизводительное ядро системы спроектировано под долгосрочное хранение и непрерывный анализ терабайтов исторической телеметрии."),
    ("📊 Аналитика для руководства (версия Pro)",
     "Автоматическая генерация глубоких статистических отчётов по состоянию оборудования для оценки эффективности работы ремонтных служб всего холдинга."),
    ("⚙ Индивидуальный интерфейс (кастомизация)",
     "Быстрая и гибкая настройка макета страниц под регламенты вашего предприятия и требования корпоративного стиля."),
    ("🏢 Унифицированный контроль",
     "Объединение разрозненных фабрик и удалённых карьеров в единую сквозную сеть мониторинга для головной компании."),
], cols=3)
picture_h(s, LOGO_COLOR, Inches(11.05), Inches(6.98), Inches(0.32))
page_no(s)

# ============================================================ SLIDE 18 — AI DIAGNOSIS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Программное обеспечение", "Профессиональный анализ и автоматическая ИИ-диагностика", title_size=21)
photo_box(s, ML, y + Inches(0.1), Inches(5.9), Inches(1.65), A + "ai_panel_left.jpg")
photo_box(s, Inches(6.75), y + Inches(0.1), Inches(5.95), Inches(1.65), A + "ai_panel_right.jpg")
tile_grid(s, ML, y + Inches(1.85), Inches(5.9), Inches(2.65), [
    ("30+ инструментов профессиональной диагностики",
     "Удовлетворяют потребности инженеров в поиске всех видов неисправностей вращающегося оборудования."),
    ("🛠 Полный набор аналитических методов",
     "Система включает спектральный анализ, анализ временных сигналов (форма волны), каскад спектров, орбиты валов и оценку трендов."),
    ("🎯 Быстрый поиск первопричин",
     "Возможность одновременного сравнения данных с нескольких измерительных точек для точного выявления несоосности, дисбаланса или ослабления опор."),
], cols=1)
tile_grid(s, Inches(6.75), y + Inches(1.85), Inches(5.95), Inches(3.5), [
    ("AI-диагностика",
     "Алгоритм автоматически даёт выводы по диагностике неисправностей и рекомендации по обслуживанию."),
    ("🤖 Умный ИИ-диагност",
     "Алгоритм автоматически сопоставляет признаки, определяет конкретный тип неисправности и выдаёт готовое ремонтное предписание."),
    ("🚨 Разделение уровней тревоги",
     "Цветовая индикация опасности (жёлтый / оранжевый / красный) по направлениям дефектов позволяет мгновенно локализовать узел, требующий внимания."),
    ("📉 Автоматический перехват аномалий",
     "Система мгновенно вырезает и сохраняет фрагмент формы волны в момент скачка вибрации для последующего детального разбора."),
], cols=1)
page_no(s)

# ============================================================ SLIDE 19 — MOBILE APP
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Программное обеспечение", "Мобильное приложение SuperCare: оперативный контроль предприятия 24/7",
              "Полный доступ к предиктивной аналитике, статусу оборудования холдинга и системе аварийных "
              "оповещений прямо со смартфона руководства и сервисных служб.",
              title_size=19)
ph_h = Inches(2.65)
ph_w = Emu(int(ph_h * (9 / 19.5)))
photo_box(s, Inches(1.8), y + Inches(0.15), ph_w, ph_h, A + "phone_left.jpg")
rect(s, Inches(5.1), y + Inches(0.6), Inches(3.1), Inches(1.75), GREY_BG)
photo_box(s, Inches(5.2), y + Inches(0.68), Inches(2.9), Inches(1.6), A + "phone_chart.jpg")
photo_box(s, Inches(9.6), y + Inches(0.15), ph_w, ph_h, A + "phone_right.jpg")
tile_grid(s, ML, y + Inches(3.0), Inches(12.1), Inches(1.85), [
    ("🟢 Контроль статуса в режиме реального времени",
     "Мгновенная проверка состояния любого агрегата (работает / остановлен / авария). Базовая техническая информация по узлам, типу оборудования и уровню важности всегда под рукой."),
    ("📈 Просмотр детальной статистики и трендов",
     "Быстрый вывод исторических графиков вибрации и трендов волновых спектров на экране телефона. Возможность оперативно оценить динамику развития дефекта без доступа к ПК."),
    ("🚨 Управление тревогами и сквозной аудит",
     "Удобный журнал обработки инцидентов. Диспетчер или главный инженер видит список необработанных аварийных сигналов с точным указанием цеха, типа неисправности и времени её фиксации."),
], cols=3)
page_no(s)

# ============================================================ SLIDE 20 — EPM
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Программное обеспечение", "Система EPM: наглядный мониторинг и экспресс-диагностика",
              "Упрощённый интерфейс платформы, созданный специально для диспетчеров и руководителей цехов. "
              "Позволяет мгновенно оценить общую ситуацию с надёжностью оборудования без необходимости "
              "работы со сложными спектральными графиками.", title_size=18)
photo_box(s, ML, y + Inches(0.1), Inches(6.4), Inches(2.6), A + "epm_dashboard.jpg")
feat_list(s, Inches(7.2), y + Inches(0.1), Inches(5.5), [
    ("📐 Наглядное 2.5D-моделирование агрегатов",
     "Реализация интуитивно понятных трёхмерных моделей оборудования и точек измерения прямо на экране."),
    ("📊 Сквозная статистика состояния парка",
     "Цветовая кодировка рисков помогает управляющему персоналу за 3 секунды понять статус надёжности цеха."),
    ("💻 Доступ без установки",
     "Программа полностью совместима с любыми современными браузерами и не требует развёртывания софта."),
], row_h=0.82)
photo_box(s, ML, y + Inches(2.85), Inches(6.15), Inches(1.7), A + "epm_machine3d.jpg",
          "Мониторинг состояния оборудования в реальном времени")
photo_box(s, Inches(6.85), y + Inches(2.85), Inches(5.85), Inches(1.7), A + "epm_waveform.jpg",
          "Профессиональные инструменты диагностики")
page_no(s)

prs.save("algoritm-full-partial.pptx")
print("saved batch 3 (slides 10-20)")


# ============================================================ SLIDE 22 — SIX SERVICES
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Сервисы поддержки", "Сервисная поддержка холдинга: экспертное сопровождение системы 24/7",
              "Мы не просто поставляем оборудование и софт — наша инженерная служба обеспечивает полный цикл "
              "экспертной поддержки для достижения максимальной эффективности и окупаемости проекта.",
              title_size=18)
photo_box(s, ML, y + Inches(0.1), Inches(5.9), Inches(1.75), A + "services_control_room.jpg",
          "Круглосуточный мониторинг в реальном времени")
photo_box(s, Inches(6.75), y + Inches(0.1), Inches(5.95), Inches(1.75), A + "services_campus.jpg",
          "Синхронизация данных в облачном диагностическом центре")
tile_grid(s, ML, y + Inches(2.0), Inches(12.1), Inches(2.55), [
    ("📡 Круглосуточный мониторинг",
     "Экспертный контроль состояния ваших агрегатов в режиме 24/7/365. Мгновенная реакция удалённого диагностического центра на любые аварийные сигналы системы."),
    ("☁ Синхронизация данных",
     "Непрерывное резервное копирование и синхронизация телеметрии в защищённом облачном диагностическом центре для исключения потери истории дефектов."),
    ("🧠 Экспертный консалтинг",
     "Постоянный доступ вашей технической службы к ведущим инженерам-диагностам для решения сложных и нестандартных инцидентов на фабрике."),
    ("📝 Ежемесячный аудит",
     "Предоставление детального отчёта о состоянии всего парка оборудования с анализом трендов развития дефектов и рекомендациями для планирования бюджетов ТОиР."),
    ("🛠 Оценка после техобслуживания",
     "Экспертный контроль качества выполненных ремонтов на месте. Проверка геометрии узлов, качества балансировки и центровки агрегатов после запуска."),
    ("👥 Постоянное обучение",
     "Регулярное повышение квалификации сотрудников вашей ИТ- и ремонтной службы для эффективной самостоятельной работы с платформой."),
], cols=3)
page_no(s)

# ============================================================ SLIDE 23 — TRAINING
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
y = page_head(s, "Сервисы поддержки", "Центр обучения персонала и авторизация экосистемы",
              "Мы помогаем заказчикам выстроить собственную независимую экосистему ТОиР на базе искусственного "
              "интеллекта и сертифицируем специалистов по высшим отраслевым стандартам.", title_size=19)
feat_list(s, ML, y + Inches(0.1), Inches(5.9), [
    ("🤖 Передача ИИ-компетенций",
     "Мы помогаем клиентам развернуть и настроить собственную внутреннюю систему O&M (эксплуатации и обслуживания), обучая встроенные ИИ-алгоритмы под задачи предприятия."),
    ("🛡 Сертифицированная платформа",
     "Внедрение решений на базе международной GL-сертифицированной платформы технического обслуживания, что гарантирует соответствие самым строгим промышленным аудитам."),
    ("📈 Сопровождение цифровой трансформации",
     "Полное управление рисками при переходе от устаревшего планового ремонта к интеллектуальному обслуживанию оборудования по его фактическому состоянию."),
], row_h=1.15)
photo_box(s, Inches(6.75), y + Inches(0.1), Inches(5.95), Inches(1.85), A + "training_classroom.jpg",
          "Обучение инженеров вибродиагностике")
photo_box(s, Inches(6.75), y + Inches(2.1), Inches(5.95), Inches(1.85), A + "training_serverroom.jpg",
          "Инфраструктура диагностического центра")
page_no(s)

# ============================================================ SLIDE 36 — CONTACTS
GREY_LIGHT = RGBColor(0xF1, 0xEF, 0xEE)
s = slide()
rect(s, 0, 0, SW, SH, GREY_LIGHT)
rect(s, 0, 0, Inches(5.87), SH, STEEL_DARK)
picture_h(s, LOGO_WHITE, Inches(0.6), Inches(2.95), Inches(0.85))
_, tf = textbox(s, Inches(0.6), Inches(4.05), Inches(4.4), Inches(0.6))
add_para(tf, "ООО «Алгоритм»", 20, WHITE, bold=True, first=True)
_, tf = textbox(s, Inches(0.6), Inches(4.65), Inches(4.3), Inches(0.8))
add_para(tf, "Комплексные решения для предиктивного обслуживания оборудования.",
         12.5, WHITE_MUTED, first=True, spacing=1.35)
_, tf = textbox(s, Inches(6.55), Inches(1.1), Inches(6.1), Inches(1.5))
add_para(tf, "Сделайте первый шаг к оптимизации ТОиР вашего предприятия", 18, INK, bold=True, first=True, spacing=1.1)
_, tf = textbox(s, Inches(6.55), Inches(1.85), Inches(5.9), Inches(0.75))
add_para(tf, "Свяжитесь с нами, чтобы согласовать проведение бесплатного экспресс-аудита критического "
             "оборудования холдинга.", 11.5, INK_MUTED, first=True, spacing=1.3)
rows = [
    ("📞 Телефон", "+7 (988) 492-81-03", "Прямая линия с инженерно-сервисным центром."),
    ("✉ E-mail", "info@algoritmkrsk.ru", "Для направления технических заданий и опросных листов."),
    ("🌐 Сайт", "algoritmkrsk.ru", None),
    ("📍 Адрес", "660075, Красноярский край, г. Красноярск, ул. Маерчака, д. 8, офис 318",
     "Собственная база поддержки и склад оборудования в СФО."),
]
ry = Inches(2.85)
for label, value, note in rows:
    _, tf = textbox(s, Inches(6.55), ry, Inches(6.0), Inches(0.3))
    add_para(tf, label, 10.5, PRIMARY, bold=True, first=True, letter_caps=True)
    val_size = 13 if label == "📍 Адрес" else 17
    _, tf = textbox(s, Inches(6.55), ry + Inches(0.3), Inches(6.0), Inches(0.55))
    add_para(tf, value, val_size, INK, bold=True, first=True, spacing=1.15)
    row_h = Inches(0.85)
    if note:
        _, tf = textbox(s, Inches(6.55), ry + Inches(0.78), Inches(6.0), Inches(0.4))
        add_para(tf, note, 10, INK_MUTED, first=True, spacing=1.2)
        row_h = Inches(1.1)
    if label != "📍 Адрес":
        rect(s, Inches(6.55), ry + row_h - Inches(0.12), Inches(6.0), Pt(0.75), RULE)
    ry = ry + row_h
page_no(s)

prs.save("algoritm-predictive-maintenance-Swiss-full.pptx")
print(f"saved FINAL: {len(prs.slides)} slides")
