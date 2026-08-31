#!/usr/bin/env python3
"""Build the Algoritm branded deck as a native PPTX (hand-mapped from the HTML deck)."""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

PRIMARY = RGBColor(0xB7, 0x2E, 0x26)
ACCENT = RGBColor(0x8C, 0x1B, 0x13)
SECONDARY = RGBColor(0x14, 0x12, 0x12)
GREY_BG = RGBColor(0xF1, 0xEF, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x14, 0x12, 0x12)
INK_MUTED = RGBColor(0x5A, 0x57, 0x57)
WHITE_MUTED = RGBColor(0xC9, 0xC4, 0xC4)
PINK = RGBColor(0xE8, 0x93, 0x8C)
RULE = RGBColor(0xD8, 0xD4, 0xD3)
RULE_DARK = RGBColor(0x3A, 0x37, 0x37)

FONT = "DM Sans"
LOGO_COLOR = "assets/algoritm-logo.png"
LOGO_WHITE = "assets/algoritm-logo-white.png"
LOGO_ASPECT = 339 / 202

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color=None, line=False):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if color is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def gradient_rect(s, l, t, w, h, c1, c2, angle=45):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    shp.line.fill.background()
    shp.fill.gradient()
    stops = shp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[-1].position = 1.0
    stops[-1].color.rgb = c2
    try:
        shp.fill.gradient_angle = angle
    except Exception:
        pass
    return shp


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_line_spacing(para, mult):
    para.line_spacing = mult


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, first=False,
             space_after=0, spacing=1.0, font=FONT, letter_caps=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text.upper() if letter_caps else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return p


def picture_h(s, path, l, t, height):
    width = Emu(int(height * LOGO_ASPECT))
    return s.shapes.add_picture(path, l, t, height=height, width=width)


def page_no(s, text, on_dark=False):
    _, tf = textbox(s, Inches(11.9), Inches(6.95), Inches(1.1), Inches(0.35))
    add_para(tf, text, 10, WHITE_MUTED if on_dark else INK_MUTED, bold=True, align=PP_ALIGN.RIGHT, first=True)


def kicker(s, l, t, w, text, color):
    _, tf = textbox(s, l, t, w, Inches(0.3))
    add_para(tf, text, 11, color, bold=True, first=True, letter_caps=True)


def connector(s, x1, y1, x2, y2, color, dashed=False, width=1.0):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.shadow.inherit = False
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dashed:
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return ln


def node_box(s, l, t, w, h, color, label, label_size=11):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    try:
        box.adjustments[0] = 0.18
    except Exception:
        pass
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    lines = label.split("\n")
    add_para(tf, lines[0], label_size, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, spacing=1.05)
    for extra in lines[1:]:
        add_para(tf, extra, label_size, WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=1.05)
    return box


def gateway_diagram(s, bx, by, bw, bh, wireless, gateway_label, protocol_label):
    """Sensors -> gateway -> server schematic (own drawing, not a factory photo)."""
    line_color = PRIMARY if wireless else RGBColor(0x6E, 0x6A, 0x69)
    bh_half = Emu(int(bh / 2))
    mid_y = by + bh_half
    sensor_x = bx + Inches(0.28)
    sensor_r = Inches(0.075)
    sensor_ys = [by + Inches(0.15), mid_y, by + bh - Inches(0.15)]

    gw_w, gw_h = Inches(1.15), Inches(0.48)
    gw_x = bx + Inches(1.0)
    gw_y = mid_y - Emu(int(gw_h / 2))

    srv_w, srv_h = Inches(1.15), Inches(0.48)
    srv_x = bx + bw - srv_w - Inches(0.15)
    srv_y = mid_y - Emu(int(srv_h / 2))

    for sy in sensor_ys:
        connector(s, sensor_x + sensor_r, sy, gw_x, mid_y, line_color, dashed=wireless, width=1.1)

    for sy in sensor_ys:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, sensor_x, sy - sensor_r, sensor_r * 2, sensor_r * 2)
        dot.shadow.inherit = False
        dot.fill.solid()
        dot.fill.fore_color.rgb = SECONDARY
        dot.line.fill.background()

    _, tf = textbox(s, sensor_x - Inches(0.35), by + bh - Inches(0.02), Inches(1.0), Inches(0.2))
    add_para(tf, "Датчики", 6.5, INK_MUTED, align=PP_ALIGN.CENTER, first=True)

    _, tf = textbox(s, gw_x - Inches(0.2), gw_y - Inches(0.2), gw_w + Inches(0.4), Inches(0.18))
    add_para(tf, protocol_label, 6, INK_MUTED, align=PP_ALIGN.CENTER, first=True)

    node_box(s, gw_x, gw_y, gw_w, gw_h, line_color if wireless else SECONDARY, gateway_label, label_size=10.5)
    _, tf = textbox(s, gw_x - Inches(0.2), gw_y + gw_h + Inches(0.02), gw_w + Inches(0.4), Inches(0.18))
    add_para(tf, "Шлюз", 6, INK_MUTED, align=PP_ALIGN.CENTER, first=True)

    connector(s, gw_x + gw_w, mid_y, srv_x, mid_y,
              PRIMARY if not wireless else SECONDARY, width=1.1)

    node_box(s, srv_x, srv_y, srv_w, srv_h, PRIMARY if not wireless else SECONDARY,
             "СЕРВЕР\nАСУ ТП", label_size=8.5)
    _, tf = textbox(s, srv_x - Inches(0.3), srv_y + srv_h + Inches(0.02), srv_w + Inches(0.6), Inches(0.18))
    add_para(tf, "Хранение и анализ", 6, INK_MUTED, align=PP_ALIGN.CENTER, first=True)


def bar_chart(s, x, y, w, h, title_lines, categories, values, bar_color, first_color, cat_color, val_color):
    _, tf = textbox(s, x, y, w, Inches(0.5))
    add_para(tf, title_lines[0], 11.5, val_color, bold=True, first=True, spacing=1.15, space_after=1)
    for extra in title_lines[1:]:
        add_para(tf, extra, 11.5, val_color, bold=True, spacing=1.15)

    title_h = Inches(0.15) * len(title_lines) + Inches(0.35)
    top = y + title_h
    y_base = y + h - Inches(0.32)
    avail_h = y_base - top - Inches(0.28)
    max_val = max(values)

    connector(s, x, y_base, x + w, y_base, RGBColor(0x8A, 0x86, 0x85), width=0.75)

    n = len(values)
    slot_w = Emu(int(w / n))
    bar_w = Emu(int(slot_w * 0.5))
    for i, (cat, val) in enumerate(zip(categories, values)):
        bar_h = Emu(int(avail_h * (val / max_val)))
        slot_x = x + Emu(slot_w * i)
        bar_x = slot_x + Emu(int((slot_w - bar_w) / 2))
        bar_y = y_base - bar_h
        color = first_color if i == 0 else bar_color
        rect(s, bar_x, bar_y, bar_w, bar_h, color)

        _, tf = textbox(s, bar_x - Inches(0.15), bar_y - Inches(0.28), bar_w + Inches(0.3), Inches(0.24))
        add_para(tf, f"{val}%", 11, color, bold=True, align=PP_ALIGN.CENTER, first=True)

        _, tf = textbox(s, slot_x, y_base + Inches(0.05), slot_w, Inches(0.3))
        add_para(tf, cat, 8, cat_color, bold=True, align=PP_ALIGN.CENTER, first=True, spacing=1.1)


# ============================================================ SLIDE 1 — COVER
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
gradient_rect(s, Inches(6.13), 0, Inches(7.2), SH, PRIMARY, ACCENT, angle=115)

picture_h(s, LOGO_COLOR, Inches(0.6), Inches(0.55), Inches(0.95))

rect(s, Inches(0.6), Inches(1.75), Inches(0.85), Inches(0.05), PRIMARY)

_, tf = textbox(s, Inches(0.6), Inches(2.0), Inches(5.1), Inches(2.7))
add_para(tf, "Предиктивное обслуживание и решения", 33, INK, bold=True, first=True, spacing=0.98)
add_para(tf, "для горнодобывающей", 33, INK, bold=True, spacing=0.98)
add_para(tf, "промышленности", 33, INK, bold=True, spacing=0.98)

_, tf = textbox(s, Inches(0.6), Inches(4.75), Inches(4.8), Inches(0.7))
add_para(tf, "ООО «Алгоритм» · Красноярск", 12.5, INK_MUTED, first=True, spacing=1.3)
add_para(tf, "Комплексный подход к решению задач", 12.5, INK_MUTED, spacing=1.3)

_, tf = textbox(s, Inches(7.8), Inches(5.2), Inches(5.1), Inches(1.2), anchor=MSO_ANCHOR.BOTTOM)
add_para(tf, "Прогнозное техническое обслуживание оборудования на базе промышленного мониторинга",
         14.5, WHITE, bold=True, align=PP_ALIGN.RIGHT, first=True, spacing=1.25)

_, tf = textbox(s, Inches(7.8), Inches(6.55), Inches(5.1), Inches(0.35))
add_para(tf, "Коммерческое предложение", 10.5, WHITE_MUTED, align=PP_ALIGN.RIGHT, first=True, letter_caps=True)

page_no(s, "01", on_dark=True)

# ============================================================ SLIDE 2 — ABOUT
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(2.55), SECONDARY)
kicker(s, Inches(0.6), Inches(0.42), Inches(5), "О компании", PRIMARY)
_, tf = textbox(s, Inches(0.6), Inches(0.8), Inches(10.5), Inches(1.5))
add_para(tf, "ООО «Алгоритм» — ваш партнер в цифровизации", 24, WHITE, bold=True, first=True, spacing=1.02)
add_para(tf, "промышленности", 24, WHITE, bold=True, spacing=1.02)

cols = [
    (Inches(0.6), Inches(3.05), None,
     "Мы предлагаем полный цикл внедрения систем прогнозного технического обслуживания: "
     "от аудита оборудования и подбора аппаратных решений до развертывания программного "
     "обеспечения и обучения персонала.", SECONDARY),
    (Inches(4.05), Inches(3.05), "Индивидуальный подход",
     "Адаптируем решения под специфику конкретного предприятия и отрасли.", PRIMARY),
    (Inches(7.15), Inches(3.05), "Экспертная поддержка",
     "Обеспечиваем консалтинг и сопровождение на всех этапах жизненного цикла системы.", PRIMARY),
    (Inches(10.25), Inches(3.05), "Мировые стандарты",
     "Опираемся на проверенную технологическую базу, успешно зарекомендовавшую себя на "
     "крупнейших промышленных объектах.", PRIMARY),
]
col_w = Inches(2.85)
for l, t, heading, body, rule_color in cols:
    rect(s, l, t, col_w, Inches(0.035), rule_color)
    y = t + Inches(0.18)
    _, tf = textbox(s, l, y, col_w, Inches(3.2))
    first = True
    if heading:
        add_para(tf, heading, 14.5, PRIMARY, bold=True, first=True, spacing=1.1, space_after=6)
        first = False
    add_para(tf, body, 12.5, INK_MUTED if heading else INK, first=first, spacing=1.35)

picture_h(s, LOGO_COLOR, Inches(0.6), Inches(6.55), Inches(0.55))
page_no(s, "02")

# ============================================================ SLIDE 3 — SCALE
s = slide()
rect(s, 0, 0, SW, SH, SECONDARY)
gradient_rect(s, 0, 0, SW, Inches(3.15), PRIMARY, ACCENT, angle=100)

kicker(s, Inches(0.6), Inches(0.42), Inches(6), "Технологическая база", RGBColor(0xF0, 0xC3, 0xBF))
_, tf = textbox(s, Inches(0.6), Inches(0.8), Inches(9.5), Inches(1.3))
add_para(tf, "Масштаб внедрения предлагаемых решений", 25, WHITE, bold=True, first=True, spacing=1.02)

picture_h(s, LOGO_WHITE, Inches(11.75), Inches(0.55), Inches(0.62))

stats = [
    ("800+", "сотрудников в командах разработки и инженерной поддержки технологической платформы"),
    ("250+", "патентов и зарегистрированных технических решений в основе платформы"),
    ("882 000+", "датчиков, установленных на промышленных объектах по всему миру"),
]
sx = [Inches(0.6), Inches(4.9), Inches(9.2)]
sw = Inches(3.9)
for (num, cap), x in zip(stats, sx):
    rect(s, x, Inches(4.55), sw, Inches(0.02), RGBColor(0x6E, 0x62, 0x60))
    _, tf = textbox(s, x, Inches(4.75), sw, Inches(1.0))
    add_para(tf, num, 40, WHITE, bold=True, first=True, spacing=0.95)
    _, tf = textbox(s, x, Inches(5.75), sw, Inches(1.0))
    add_para(tf, cap, 12.5, WHITE_MUTED, first=True, spacing=1.3)

_, tf = textbox(s, Inches(0.6), Inches(6.95), Inches(11.5), Inches(0.4))
add_para(tf, "Показатели характеризуют технологическую базу, на которую опирается ООО «Алгоритм» "
             "при реализации проектов в РФ.", 10.5, RGBColor(0x8A, 0x86, 0x85), first=True)

# ============================================================ SLIDE 4 — BASIS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, Inches(5.07), SH, GREY_BG)

kicker(s, Inches(0.6), Inches(1.4), Inches(4), "Основание для выбора", PRIMARY)
_, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(4.0), Inches(1.9))
add_para(tf, "Проверенная технология в руках", 20, INK, bold=True, first=True, spacing=1.05)
add_para(tf, "российского интегратора", 20, INK, bold=True, spacing=1.05)

_, tf = textbox(s, Inches(0.6), Inches(3.2), Inches(4.0), Inches(1.0))
add_para(tf, "ООО «Алгоритм» выступает единым центром ответственности: поставка, внедрение, "
             "поддержка и обучение на территории РФ.", 12.5, INK_MUTED, first=True, spacing=1.35)

rect(s, Inches(0.6), Inches(4.55), Inches(0.035), Inches(0.85), PRIMARY)
_, tf = textbox(s, Inches(0.8), Inches(4.55), Inches(3.8), Inches(0.9))
add_para(tf, "Слайды 3–4 исходной презентации переработаны: акцент смещён с юрлица "
             "завода-изготовителя на масштаб внедрения технологии.", 10.5, ACCENT, first=True, spacing=1.3)

steps = [
    ("01", "Аудит и обследование", "Определяем критичное оборудование, точки контроля и целевые показатели надёжности."),
    ("02", "Подбор аппаратных решений", "Формируем конфигурацию датчиков и шлюзов под условия конкретной площадки."),
    ("03", "Развёртывание ПО", "Пусконаладка системы мониторинга и интеграция с существующим контуром предприятия."),
    ("04", "Обучение и сопровождение", "Передаём компетенции персоналу заказчика и обеспечиваем поддержку на всём жизненном цикле."),
]
ry = Inches(1.15)
row_h = Inches(1.42)
for idx, heading, body in steps:
    _, tf = textbox(s, Inches(5.65), ry, Inches(0.9), Inches(0.7))
    add_para(tf, idx, 22, PRIMARY, bold=True, first=True)
    _, tf = textbox(s, Inches(6.6), ry, Inches(6.1), Inches(1.15))
    add_para(tf, heading, 15, INK, bold=True, first=True, spacing=1.1, space_after=4)
    add_para(tf, body, 12.5, INK_MUTED, spacing=1.3)
    if idx != "04":
        rect(s, Inches(5.65), ry + row_h - Inches(0.12), Inches(7.05), Emu(9525), RULE)
    ry = ry + row_h

page_no(s, "04")

# ============================================================ SLIDE 5 — SCENARIOS
s = slide()
rect(s, 0, 0, SW, SH, GREY_BG)
kicker(s, Inches(0.6), Inches(0.42), Inches(5), "Сценарии применения", PRIMARY)
_, tf = textbox(s, Inches(0.6), Inches(0.8), Inches(9.5), Inches(1.2))
add_para(tf, "Беспроводные и проводные схемы сбора данных", 22, INK, bold=True, first=True, spacing=1.02)

picture_h(s, LOGO_COLOR, Inches(11.75), Inches(0.5), Inches(0.62))

cards = [
    (Inches(0.6), "Беспроводной шлюз", "RH570", PRIMARY, True, "LoRaWAN · Wi-Fi",
     "Развёртывание без прокладки кабельных трасс — для распределённого и труднодоступного оборудования."),
    (Inches(6.75), "Проводной шлюз", "RH2000", SECONDARY, False, "RS-485 · Ethernet",
     "Стационарная схема для ответственных агрегатов с высокой частотой опроса и постоянным питанием."),
]
card_w = Inches(6.0)
for x, tag, title, rule_color, wireless, protocol, body in cards:
    rect(s, x, Inches(2.15), card_w, Inches(0.045), rule_color)
    rect(s, x, Inches(2.2), card_w, Inches(2.75), WHITE)
    _, tf = textbox(s, x + Inches(0.35), Inches(2.42), card_w - Inches(0.7), Inches(0.3))
    add_para(tf, tag, 10.5, RGBColor(0x8A, 0x86, 0x85), bold=True, first=True, letter_caps=True)
    _, tf = textbox(s, x + Inches(0.35), Inches(2.72), card_w - Inches(0.7), Inches(0.5))
    add_para(tf, title, 22, INK, bold=True, first=True)

    diagram_w = card_w - Inches(0.7)
    diagram_h = Inches(1.2)
    diagram_y = Inches(3.3)
    rect(s, x + Inches(0.35), diagram_y, diagram_w, diagram_h, GREY_BG)
    gateway_diagram(s, x + Inches(0.35), diagram_y + Inches(0.15), diagram_w, diagram_h - Inches(0.3),
                     wireless, title, protocol)

    _, tf = textbox(s, x + Inches(0.35), Inches(4.55), card_w - Inches(0.7), Inches(0.42))
    add_para(tf, body, 12, INK_MUTED, first=True, spacing=1.25)

page_no(s, "12")

# ============================================================ SLIDE 6 — SERVICES
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
gradient_rect(s, Inches(10.13), 0, Inches(3.2), SH, PRIMARY, ACCENT, angle=160)

kicker(s, Inches(0.6), Inches(0.42), Inches(5), "Сервисы поддержки", PRIMARY)
_, tf = textbox(s, Inches(0.6), Inches(0.8), Inches(9.0), Inches(1.6))
add_para(tf, "Услуги, которые ООО «Алгоритм» оказывает своим", 22, INK, bold=True, first=True, spacing=1.05)
add_para(tf, "клиентам в РФ", 22, INK, bold=True, spacing=1.05)

services = [
    ("Круглосуточный мониторинг", "Непрерывный контроль состояния оборудования и оповещение об отклонениях."),
    ("Обучение персонала", "Подготовка специалистов заказчика к самостоятельной работе с системой."),
    ("Консалтинг и сопровождение", "Методическая поддержка на всех этапах жизненного цикла системы."),
]
sy = Inches(2.7)
for heading, body in services:
    rect(s, Inches(0.6), sy, Inches(0.035), Inches(0.75), PRIMARY)
    _, tf = textbox(s, Inches(0.85), sy, Inches(8.3), Inches(0.8))
    add_para(tf, heading, 15, INK, bold=True, first=True, spacing=1.1, space_after=3)
    add_para(tf, body, 12.5, INK_MUTED, spacing=1.3)
    sy = sy + Inches(0.95)

badge_cx, badge_cy, badge_r = Inches(10.8), Inches(5.15), Inches(0.28)
badge = s.shapes.add_shape(MSO_SHAPE.OVAL, badge_cx - badge_r, badge_cy - badge_r, badge_r * 2, badge_r * 2)
badge.shadow.inherit = False
badge.fill.background()
badge.line.color.rgb = WHITE
badge.line.width = Pt(1.5)
connector(s, badge_cx - Inches(0.12), badge_cy, badge_cx - Inches(0.02), badge_cy + Inches(0.1), WHITE, width=2.2)
connector(s, badge_cx - Inches(0.02), badge_cy + Inches(0.1), badge_cx + Inches(0.15), badge_cy - Inches(0.12), WHITE, width=2.2)

_, tf = textbox(s, Inches(10.45), Inches(5.55), Inches(2.6), Inches(1.6), anchor=MSO_ANCHOR.BOTTOM)
add_para(tf, "Сертификация и соответствие", 15.5, WHITE, bold=True, first=True, spacing=1.1, space_after=8)
add_para(tf, "Подтверждающие документы производителя приводятся справочно, в приложении к предложению.",
         11.5, RGBColor(0xF3, 0xD8, 0xD5), spacing=1.3)

page_no(s, "21")

# ============================================================ SLIDE 7 — CASES
s = slide()
rect(s, 0, 0, SW, SH, SECONDARY)
diamond = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(-3.0), Inches(9.5), Inches(9.5))
diamond.rotation = 45
diamond.shadow.inherit = False
diamond.line.fill.background()
diamond.fill.gradient()
dstops = diamond.fill.gradient_stops
dstops[0].position = 0.0
dstops[0].color.rgb = PRIMARY
dstops[-1].position = 1.0
dstops[-1].color.rgb = SECONDARY
try:
    diamond.fill.gradient_angle = 135
except Exception:
    pass

kicker(s, Inches(0.6), Inches(2.1), Inches(4), "Раздел", RGBColor(0xF0, 0xC3, 0xBF))
_, tf = textbox(s, Inches(0.6), Inches(2.5), Inches(7.0), Inches(1.9))
add_para(tf, "Международный опыт применения технологии на", 24, WHITE, bold=True, first=True, spacing=1.02)
add_para(tf, "крупнейших промышленных объектах", 24, WHITE, bold=True, spacing=1.02)

_, tf = textbox(s, Inches(0.6), Inches(4.55), Inches(6.6), Inches(0.9))
add_para(tf, "Раздел объединяет проекты внедрения на предприятиях горнодобывающей, "
             "металлургической и энергетической отраслей.", 13, WHITE_MUTED, first=True, spacing=1.35)

tags = ["Горная добыча", "Металлургия", "Энергетика", "Цементное производство"]
tx = Inches(0.6)
for t in tags:
    tw = Inches(0.35 + 0.11 * len(t))
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, Inches(5.55), tw, Inches(0.42))
    box.shadow.inherit = False
    box.fill.background()
    box.line.color.rgb = RGBColor(0x5A, 0x56, 0x56)
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    add_para(tf, t, 11.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tx = tx + tw + Inches(0.15)

rect(s, Inches(0.6), Inches(6.15), Inches(0.035), Inches(0.6), PINK)
_, tf = textbox(s, Inches(0.82), Inches(6.15), Inches(5.6), Inches(0.7))
add_para(tf, "Слайды 27–33 исходной презентации собираются под этим заголовком. "
             "На слайде 29 фотография убирается — остаётся только график.", 10.5, PINK, first=True, spacing=1.3)

bar_chart(
    s, Inches(9.05), Inches(1.2), Inches(3.65), Inches(2.55),
    ["Снижение внеплановых простоев,", "по направлениям внедрения"],
    ["Добыча", "Металлургия", "Энергетика", "Цемент"],
    [34, 28, 22, 19],
    bar_color=RGBColor(0xF3, 0xD6, 0xD3), first_color=WHITE,
    cat_color=RGBColor(0xE8, 0xD8, 0xD6), val_color=WHITE,
)
_, tf = textbox(s, Inches(9.05), Inches(3.95), Inches(3.65), Inches(0.6))
add_para(tf, "Иллюстративные данные — приводятся как референс и уточняются по факту внедрения.",
         9, RGBColor(0xC9, 0xAE, 0xAC), first=True, spacing=1.3)

picture_h(s, LOGO_WHITE, Inches(11.75), Inches(0.5), Inches(0.62))

# ============================================================ SLIDE 8 — CONTACTS
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, Inches(5.87), SH, SECONDARY)

picture_h(s, LOGO_WHITE, Inches(0.6), Inches(2.95), Inches(0.85))
_, tf = textbox(s, Inches(0.6), Inches(4.05), Inches(4.4), Inches(0.6))
add_para(tf, "ООО «Алгоритм»", 20, WHITE, bold=True, first=True)
_, tf = textbox(s, Inches(0.6), Inches(4.65), Inches(4.3), Inches(0.8))
add_para(tf, "Комплексные решения для предиктивного обслуживания оборудования.",
         12.5, WHITE_MUTED, first=True, spacing=1.35)

rows = [
    ("Телефон", "+7 (988) 492-81-03"),
    ("E-mail", "info@algoritmkrsk.ru"),
    ("Сайт", "https://algoritmkrsk.ru/"),
    ("Адрес", "660075, Красноярский край, г. Красноярск,\nул. Маерчака, д. 8, офис 318"),
]
ry = Inches(2.15)
for label, value in rows:
    _, tf = textbox(s, Inches(6.55), ry, Inches(6.0), Inches(0.3))
    add_para(tf, label, 10.5, PRIMARY, bold=True, first=True, letter_caps=True)
    _, tf = textbox(s, Inches(6.55), ry + Inches(0.32), Inches(6.0), Inches(0.9))
    lines = value.split("\n")
    add_para(tf, lines[0], 17, INK, bold=True, first=True, spacing=1.15)
    for extra in lines[1:]:
        add_para(tf, extra, 17, INK, bold=True, spacing=1.15)
    if label != "Адрес":
        rect(s, Inches(6.55), ry + Inches(1.15), Inches(6.0), Emu(9525), RULE)
    ry = ry + Inches(1.3)

page_no(s, "33")

prs.save("algoritm-predictive-maintenance-Swiss.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst) if False else len(prs.slides._sldIdLst), "slides")
