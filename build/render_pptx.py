# -*- coding: utf-8 -*-
"""Независимый рендер сохранённого PPTX: читает файл и рисует его содержимое.
Проверяет, что реально лежит в PPTX (координаты, цвета, текст, картинки)."""
import io, math, pathlib, re, sys
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

ROOT = pathlib.Path(".").resolve()
SRC = ROOT / "restavraciya-obektov-Luxury.pptx"
OUT = ROOT / "build/pptx_shots"; OUT.mkdir(parents=True, exist_ok=True)
for f in OUT.glob("*.png"): f.unlink()

W, H = 1920, 1080
prs = Presentation(SRC)
SX = W / prs.slide_width
SY = H / prs.slide_height

FONTS = {
    ("Cormorant Garamond Light", False): "/root/.fonts/CormorantGaramond-Light.ttf",
    ("Noto Sans", False): "/root/.fonts/NotoSans-Regular.ttf",
    ("Noto Sans", True): "/root/.fonts/NotoSans-SemiBold.ttf",
}
def font_for(name, bold, px):
    path = FONTS.get((name, bold)) or FONTS.get((name, False)) or FONTS[("Noto Sans", False)]
    return ImageFont.truetype(path, max(1, int(round(px))))

def px(v, s): return v * s

def solid_of(spPr):
    sf = spPr.find(qn("a:solidFill"))
    if sf is None: return None
    c = sf.find(qn("a:srgbClr"))
    if c is None: return None
    rgbv = tuple(int(c.get("val")[i:i+2], 16) for i in (0, 2, 4))
    a = c.find(qn("a:alpha"))
    return rgbv, (int(a.get("val")) / 100000 if a is not None else 1.0)

def grad_of(spPr):
    gf = spPr.find(qn("a:gradFill"))
    if gf is None: return None
    stops = []
    for gs in gf.find(qn("a:gsLst")):
        c = gs.find(qn("a:srgbClr"))
        a = c.find(qn("a:alpha"))
        stops.append((int(gs.get("pos")) / 100000,
                      tuple(int(c.get("val")[i:i+2], 16) for i in (0, 2, 4)),
                      int(a.get("val")) / 100000 if a is not None else 1.0))
    ang = gf.find(qn("a:lin"))
    return stops, (int(ang.get("ang")) / 60000 if ang is not None else 0)

def blend(base, col, a):
    return Image.blend(base, Image.new("RGB", base.size, col), a)

for si, slide in enumerate(prs.slides):
    img = Image.new("RGB", (W, H), (14, 13, 12))
    for sh in slide.shapes:
        x, y = px(sh.left, SX), px(sh.top, SY)
        w, h = px(sh.width, SX), px(sh.height, SY)
        if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
            im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
            im = im.resize((max(1, int(round(w))), max(1, int(round(h)))), Image.LANCZOS)
            img.paste(im, (int(round(x)), int(round(y))))
            continue
        spPr = sh._element.spPr if hasattr(sh._element, "spPr") else None
        if spPr is not None:
            g = grad_of(spPr)
            if g:
                stops, ang = g
                box = img.crop((int(x), int(y), int(x + w), int(y + h)))
                ov = Image.new("RGB", box.size); al = Image.new("L", box.size)
                dx, dy = math.cos(math.radians(ang)), math.sin(math.radians(ang))
                dw = ImageDraw.Draw(ov); da = ImageDraw.Draw(al)
                n = 120
                for i in range(n):
                    t = i / (n - 1)
                    for k in range(len(stops) - 1):
                        if stops[k][0] <= t <= stops[k + 1][0]:
                            f = (t - stops[k][0]) / max(1e-6, stops[k + 1][0] - stops[k][0])
                            col = tuple(int(stops[k][1][j] + (stops[k+1][1][j]-stops[k][1][j]) * f) for j in range(3))
                            av = stops[k][2] + (stops[k+1][2] - stops[k][2]) * f
                            break
                    px0 = t * box.size[0]
                    dw.rectangle([px0, 0, px0 + box.size[0]/n + 1, box.size[1]], fill=col)
                    da.rectangle([px0, 0, px0 + box.size[0]/n + 1, box.size[1]], fill=int(av*255))
                box.paste(ov, (0, 0), al)
                img.paste(box, (int(x), int(y)))
                continue
            s = solid_of(spPr)
            if s:
                col, a = s
                if a >= 0.999:
                    ImageDraw.Draw(img).rectangle([x, y, x + w - 1, y + h - 1], fill=col)
                else:
                    box = img.crop((int(x), int(y), int(x+w), int(y+h)))
                    img.paste(blend(box, col, a), (int(x), int(y)))
            ln = spPr.find(qn("a:ln"))
            if ln is not None and ln.find(qn("a:solidFill")) is not None:
                c = ln.find(qn("a:solidFill")).find(qn("a:srgbClr"))
                col = tuple(int(c.get("val")[i:i+2], 16) for i in (0, 2, 4))
                al = c.find(qn("a:alpha"))
                av = int(al.get("val"))/100000 if al is not None else 1.0
                d = ImageDraw.Draw(img, "RGBA")
                d.rectangle([x, y, x + w - 1, y + h - 1], outline=col + (int(av*255),), width=1)
        # текст
        if not sh.has_text_frame: continue
        tf = sh.text_frame
        cy = y
        for p in tf.paragraphs:
            lnSpc = p._p.find(qn("a:pPr"))
            lh = None
            if lnSpc is not None:
                sp = lnSpc.find(qn("a:lnSpc"))
                if sp is not None:
                    pts = sp.find(qn("a:spcPts"))
                    if pts is not None: lh = int(pts.get("val")) / 100 * 2   # pt -> px (1px=0.5pt)
            cx = x + 3
            maxsz = 0
            for r in p.runs:
                rPr = r._r.find(qn("a:rPr"))
                sz = int(rPr.get("sz")) / 100 * 2 if rPr is not None and rPr.get("sz") else 36
                spc = int(rPr.get("spc")) / 100 * 2 if rPr is not None and rPr.get("spc") else 0
                b = rPr is not None and rPr.get("b") == "1"
                lat = rPr.find(qn("a:latin")) if rPr is not None else None
                fam = lat.get("typeface") if lat is not None else "Noto Sans"
                col = (0, 0, 0)
                sf = rPr.find(qn("a:solidFill")) if rPr is not None else None
                if sf is not None:
                    c = sf.find(qn("a:srgbClr"))
                    if c is not None: col = tuple(int(c.get("val")[i:i+2], 16) for i in (0, 2, 4))
                ft = font_for(fam, b, sz)
                maxsz = max(maxsz, sz)
                d = ImageDraw.Draw(img)
                asc = ft.getmetrics()[0]
                base = cy + ((lh or sz * 1.2) - sz) / 2
                for ch in r.text:
                    d.text((cx, base), ch, font=ft, fill=col)
                    cx += d.textlength(ch, font=ft) + spc
            cy += lh or (maxsz * 1.2)
    img.save(OUT / f"p{si:02d}.png")
print("rendered", len(prs.slides.__iter__.__self__._sldIdLst), "slides ->", OUT)
