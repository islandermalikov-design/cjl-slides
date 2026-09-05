# -*- coding: utf-8 -*-
"""Сборка редактируемого PPTX по геометрии, снятой с HTML-версии."""
import json, pathlib, re, sys, copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, "build")
import deck_spec as S

ROOT = pathlib.Path(".").resolve()
LAYOUT = json.loads((ROOT / "build/layout.json").read_text(encoding="utf-8"))
OUT = ROOT / "restavraciya-obektov-Luxury.pptx"

PX = 6350                      # EMU на 1 px эталонного слайда 1920x1080
PT = 0.5                       # pt на 1 px
SLIDE_W, SLIDE_H = Emu(1920 * PX), Emu(1080 * PX)

FONT_SERIF = "Cormorant Garamond Light"   # начертание Light — подпись стиля
FONT_SANS = "Noto Sans"                   # кириллическое начертание (DM Sans без кириллицы)

# подписи слотов для панели выделения PowerPoint
SLOT_LABEL = dict(S.slots())


def rgb(css):
    m = re.findall(r"[\d.]+", css)
    if css.startswith("#"):
        return RGBColor.from_string(css.lstrip("#").upper())
    return RGBColor(int(float(m[0])), int(float(m[1])), int(float(m[2])))


def alpha_of(css):
    m = re.findall(r"[\d.]+", css)
    return float(m[3]) if css.startswith("rgba") and len(m) >= 4 else 1.0


def E(px_):
    return Emu(int(round(px_ * PX)))


def set_name(shape, name):
    shape._element.nvSpPr.cNvPr.set("name", name) if hasattr(shape._element, "nvSpPr") \
        else shape._element._nvXxPr.cNvPr.set("name", name)


def name_shape(shape, name, descr=None):
    el = shape._element
    cNvPr = el.find(".//" + qn("p:cNvPr"))
    cNvPr.set("name", name)
    if descr:
        cNvPr.set("descr", descr)


def rect(slide, x, y, w, h, color, alpha=1.0, name="Фигура"):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(max(w, 0.4)), E(max(h, 0.4)))
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if alpha < 1.0:
        srgb = sh.fill.fore_color._xFill.find(qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(alpha * 100000)))
    sh.text_frame.text = ""
    name_shape(sh, name)
    return sh


def scrim(slide, r, name="ЗАТЕМНЕНИЕ (SCRIM)"):
    """Линейный градиент затемнения обложки — как в HTML."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(r["x"]), E(r["y"]), E(r["w"]), E(r["h"]))
    sh.line.fill.background(); sh.shadow.inherit = False
    spPr = sh._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="141210"><a:alpha val="93000"/></a:srgbClr></a:gs>'
        '<a:gs pos="46000"><a:srgbClr val="141210"><a:alpha val="80000"/></a:srgbClr></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="141210"><a:alpha val="62000"/></a:srgbClr></a:gs>'
        '</a:gsLst><a:lin ang="1050000" scaled="0"/></a:gradFill>')
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(spPr), etree.fromstring(xml))
    name_shape(sh, name)
    return sh


def picture(slide, r, path, name, descr):
    pic = slide.shapes.add_picture(str(path), E(r["x"]), E(r["y"]), E(r["w"]), E(r["h"]))
    name_shape(pic, name, descr)
    return pic


def textbox(slide, t):
    pad = 6                                    # компенсация внутренних полей PowerPoint
    box = slide.shapes.add_textbox(E(t["x"] - pad / 2), E(t["y"]),
                                   E(t["w"] + pad), E(t["h"] + 8))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)

    txt = t["text"]
    if t.get("transform") == "uppercase":
        txt = txt.upper()
    lines = txt.split("\n")

    serif = "Cormorant" in t["family"]
    family = FONT_SERIF if serif else FONT_SANS
    size = Pt(round(t["size"] * PT, 1))
    color = rgb(t["color"])
    ls_px = 0.0 if t["ls"] in ("normal", "") else float(re.findall(r"-?[\d.]+", t["ls"])[0])
    lh = t["lh"]
    lh_pt = round(float(re.findall(r"[\d.]+", lh)[0]) * PT, 1) if lh not in ("normal",) else None

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # два прогона для «01/06»
        parts = ([(line.split("/")[0], color), ("/" + line.split("/")[1], rgb("rgb(122,113,104)"))]
                 if "idx" in str(t["tag"]) and "/" in line else [(line, color)])
        for s_, c_ in parts:
            run = p.add_run(); run.text = s_
            f = run.font
            f.size = size; f.name = family; f.color.rgb = c_
            f.bold = int(t["weight"]) >= 600
            rPr = run._r.get_or_add_rPr()
            if ls_px:
                rPr.set("spc", str(int(round(ls_px * PT * 100))))
            # a:latin ставится через font.name; порядок дочерних узлов rPr менять нельзя
        if lh_pt:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            ln = etree.SubElement(pPr, qn("a:lnSpc"))
            pts = etree.SubElement(ln, qn("a:spcPts")); pts.set("val", str(int(lh_pt * 100)))
            pPr.insert(0, ln)
    return box


prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]

for s in LAYOUT:
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, 1920, 1080, rgb(s["bg"]), name="ФОН")

    for p_ in s["panels"]:
        rect(sl, p_["x"], p_["y"], p_["w"], p_["h"], rgb(p_["bg"]), name="ПАНЕЛЬ")

    for f in s["frames"]:
        slot = f["slot"]
        if f["ph"]:
            path = ROOT / f"build/ph/{slot}.png"
            picture(sl, f, path, f"ФОТО-ПЛЕЙСХОЛДЕР — {slot}",
                    f"Место под будущую фотографию: {SLOT_LABEL.get(slot, slot)}. "
                    f"Заменить: правый клик по изображению — «Изменить рисунок».")
        else:
            picture(sl, f, ROOT / S.PHOTOS[slot], f"ФОТО — {slot}",
                    f"Фотография: {SLOT_LABEL.get(slot, slot)}. "
                    f"Заменить: правый клик — «Изменить рисунок».")

    if s["scrim"]:
        scrim(sl, s["scrim"])

    for r in s["rules"]:
        col = r["color"] or "#DCD3C6"
        rect(sl, r["x"], r["y"], r["w"], r["h"] if r["kind"] != "tile-border" else r["h"],
             rgb(col), alpha_of(col),
             name="РАМКА ПЛИТКИ" if r["kind"] == "tile-border" else "ЛИНИЯ") \
            if r["kind"] != "tile-border" else None
        if r["kind"] == "tile-border":
            sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(r["x"]), E(r["y"]), E(r["w"]), E(r["h"]))
            sh.fill.background(); sh.shadow.inherit = False
            sh.line.color.rgb = rgb(col); sh.line.width = Pt(0.5)
            lnEl = sh._element.spPr.find(qn("a:ln"))
            sc = lnEl.find(qn("a:solidFill")).find(qn("a:srgbClr"))
            a = etree.SubElement(sc, qn("a:alpha")); a.set("val", str(int(alpha_of(col) * 100000)))
            name_shape(sh, "РАМКА ПЛИТКИ")

    for t in s["texts"]:
        if "ph-cap" in str(t["tag"]):
            continue                       # подпись уже нарисована на placeholder-изображении
        textbox(sl, t)

prs.save(OUT)
print("OK", OUT, round(OUT.stat().st_size / 1024 / 1024, 2), "MB",
      "слайдов:", len(prs.slides.__iter__.__self__._sldIdLst))
